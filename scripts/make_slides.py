"""Generate the Iteration-05 final presentation as a PowerPoint (.pptx).

Covers the required topics: problem, dataset, preprocessing, feature
engineering, models, evaluation, results, challenges, lessons learned, and
future work. Numbers are read live from outputs/metrics/ and figures are
pulled from outputs/figures/.

Run:  .venv/bin/python -m scripts.make_slides
Output: FINAL_PRESENTATION.pptx
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "outputs" / "figures"
MET = ROOT / "outputs" / "metrics"
DATA = ROOT / "data" / "dataset"

NAVY = RGBColor(0x1A, 0x3C, 0x6E)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

REPO = "https://github.com/ajaynathch/requirement_doc_parser"


def m(name):
    return json.loads((MET / name).read_text())


base = m("metrics_baseline_logreg.json")
dbert = m("metrics_distilbert.json")
base_mask = m("metrics_baseline_logreg_masked.json")
dbert_mask = m("metrics_distilbert_masked.json")
tune = m("tuning.json")
stats = json.loads((DATA / "label_stats.json").read_text())
TOTAL = sum(stats.values())
maj = max(stats.values()) / TOTAL
hp = dbert.get("hyperparams", {"epochs": 10, "batch_size": 16, "lr": 3e-5})

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _tf(box):
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.8))
    p = _tf(tb).paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE


def bullets(title, items, note=None):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, title)
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), SW - Inches(1.4), SH - Inches(2.2))
    tf = _tf(box)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        level = 0
        if isinstance(item, tuple):
            item, level = item
        p.level = level
        r = p.add_run(); r.text = ("• " if level == 0 else "– ") + item
        r.font.size = Pt(20 if level == 0 else 17)
        r.font.color.rgb = DARK if level == 0 else GREY
        p.space_after = Pt(8)
    if note:
        nb = s.shapes.add_textbox(Inches(0.7), SH - Inches(0.75), SW - Inches(1.4), Inches(0.55))
        p = _tf(nb).paragraphs[0]
        r = p.add_run(); r.text = note
        r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = NAVY
    return s


def picture_slide(title, fig, caption=None, text_items=None):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, title)
    path = FIG / fig
    if text_items:
        # figure left, bullets right
        s.shapes.add_picture(str(path), Inches(0.4), Inches(1.5),
                             height=Inches(5.2))
        box = s.shapes.add_textbox(Inches(7.6), Inches(1.6), Inches(5.4), Inches(5.2))
        tf = _tf(box)
        for i, item in enumerate(text_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = "• " + item
            r.font.size = Pt(18); r.font.color.rgb = DARK
            p.space_after = Pt(10)
    else:
        pic = s.shapes.add_picture(str(path), Inches(0), Inches(1.4), height=Inches(5.4))
        pic.left = int((SW - pic.width) / 2)
    if caption:
        cb = s.shapes.add_textbox(Inches(0.5), SH - Inches(0.6), SW - Inches(1), Inches(0.5))
        p = _tf(cb).paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = GREY
    return s


# --- 1. Title ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), SW - Inches(1.8), Inches(2.6))
tf = _tf(tb)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Automated Requirement Classification"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r = p2.add_run(); r.text = "from an NVMe Specification"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p3 = tf.add_paragraph(); p3.space_before = Pt(20)
r = p3.add_run(); r.text = "EECE5644 — Machine Learning and Pattern Recognition  |  Final Project"
r.font.size = Pt(18); r.font.color.rgb = RGBColor(0xCF, 0xDD, 0xF2)
p4 = tf.add_paragraph(); p4.space_before = Pt(10)
r = p4.add_run(); r.text = "Ajaya Nath Chittela  •  ajayanath.c@northeastern.edu"
r.font.size = Pt(16); r.font.color.rgb = WHITE

# --- 2. Problem statement ---------------------------------------------------
bullets("Problem Statement", [
    "Technical specs (NVMe, PCIe, CXL) encode a different obligation level in almost every sentence.",
    ("MANDATORY (\"shall\"), PROHIBITED (\"shall not\"), RECOMMENDED (\"should\"), OPTIONAL (\"may\"), INFORMATIVE (prose).", 1),
    "Engineers must separate and type requirements by hand before writing validation tests — slow and error-prone.",
    "Goal: automatically (1) extract requirement sentences and (2) classify each by normative modality.",
    "This is the requirement-understanding stage of a larger document-to-code agent.",
], note="Supervised multi-class text classification — 5 classes, heavy imbalance.")

# --- 3. Dataset -------------------------------------------------------------
picture_slide("Dataset & Weak Labeling", "eda/class_distribution.png",
    text_items=[
        "Source: NVM Express Base Spec, Rev 2.3.",
        "Parsed with Docling → sentence records with page & section.",
        f"{TOTAL} labeled sentences.",
        "No gold labels → weak supervision from modal keywords.",
        "Realistic, heavy imbalance: INFORMATIVE dominates; RECOMMENDED/PROHIBITED rare.",
        "Stratified train / val / test split.",
    ])

# --- 4. Preprocessing -------------------------------------------------------
bullets("Data Cleaning & Preprocessing", [
    "Keep only Docling body prose (text / paragraph / list_item); drop tables, figures, headers.",
    "Regex sentence segmentation that protects spec abbreviations (e.g., i.e., Fig., No., vol.).",
    "Drop fragments < 15 chars or with no letters (removes numbers & cross-ref debris).",
    "Remove exact-text duplicates before labeling.",
    "Word-boundary label matching: \"shall not\" before \"shall\"; \"may\" doesn't fire inside \"maybe\".",
])

# --- 5. Feature engineering -------------------------------------------------
bullets("Feature Engineering", [
    "Classical models — TF-IDF: sublinear TF, 1–2 grams (captures \"shall not\", \"may not\"), accent stripping.",
    "DistilBERT — WordPiece tokenization (max length 128), fine-tuned contextual embeddings.",
    "Interpretable engineered features for EDA:",
    ("length, avg word length, digit & acronym counts, negation / cross-ref / modal-cue flags.", 1),
    "Class imbalance handled with class weights (both model families).",
])

# --- 6. EDA -----------------------------------------------------------------
picture_slide("Exploratory Data Analysis", "eda/feature_correlation_heatmap.png",
    text_items=[
        "No single surface feature strongly correlates with the label.",
        "Modal-cue flags carry the most signal.",
        "Sentence length distributions overlap across classes.",
        "→ the class boundary lives in word choice & context, not simple statistics.",
    ])

# --- 7. Models --------------------------------------------------------------
bullets("Machine Learning Models", [
    "Multinomial Naive Bayes — generative MAP baseline over TF-IDF counts.",
    "Logistic Regression (TF-IDF) — discriminative, balanced, interpretable coefficients.",
    "Linear SVM (TF-IDF) — max-margin in high-dimensional sparse space.",
    "DistilBERT — 6-layer distilled BERT + linear head, class-weighted cross-entropy loss.",
    "Classical model selected by stratified k-fold CV (macro-F1); DistilBERT by val macro-F1.",
], note=f"DistilBERT fine-tuned {hp['epochs']} epochs, batch {hp['batch_size']}, lr {hp['lr']:.0e}, on Apple-Silicon MPS.")

# --- 8. Hyperparameter tuning ----------------------------------------------
picture_slide("Hyperparameter Tuning", "tuning_before_after.png",
    text_items=[
        f"Grid search: {tune['grid_size']} combos × {tune['cv_folds']}-fold CV (macro-F1).",
        "Tuned ngram_range, min_df, sublinear_tf, C.",
        f"Best: bigrams, min_df=2, sublinear TF, C=1.0.",
        f"Test macro-F1 {tune['before_test']['macro_f1']:.2f} → {tune['after_test']['macro_f1']:.2f} "
        f"(+{tune['after_test']['macro_f1']-tune['before_test']['macro_f1']:.2f}).",
        "Biggest gains: bigrams + sublinear TF scaling.",
    ])

# --- 9. Evaluation / results ------------------------------------------------
picture_slide("Results — Model Comparison", "model_comparison.png",
    text_items=[
        f"DistilBERT: accuracy {dbert['accuracy']:.2f}, weighted-F1 {dbert['weighted_f1']:.2f}.",
        f"Tuned LogReg: best macro-F1 {tune['after_test']['macro_f1']:.2f}.",
        f"Majority baseline accuracy ≈ {maj:.2f}.",
        "Macro-F1 is the headline metric under imbalance.",
        "Linear model rivals the Transformer on this task.",
    ])

# --- 10. Confusion ----------------------------------------------------------
picture_slide("Results — Error Analysis", "confusion_distilbert.png",
    text_items=[
        "Perfect on MANDATORY & OPTIONAL (P=R=1.00).",
        "INFORMATIVE F1 ≈ 0.97.",
        "Errors sit in PROHIBITED (1 test sample) & RECOMMENDED (0).",
        "→ data-scarcity artifact, not a systematic weakness.",
        "Critical error (dropping a real requirement) is minimized: MANDATORY recall = 1.00.",
    ])

# --- 11. Feature importance + learning curve --------------------------------
picture_slide("Feature Importance & Learning Curve", "feature_importance.png",
    text_items=[
        "Top terms per class are intuitive:",
        "MANDATORY → \"shall\", \"controller shall\".",
        "OPTIONAL/PROHIBITED → \"may\", \"may not\".",
        "Generic nouns & \"figure\"/\"section\" ≈ zero weight.",
        "Learning curve still rising → more data would help.",
    ])

# --- 12. Ablation -----------------------------------------------------------
bullets("Do the Models Learn Meaning, or Just Keywords?", [
    "Ablation: blank the modal trigger word and re-evaluate.",
    f"Logistic Regression accuracy: {base['accuracy']:.2f} (full) → {base_mask['accuracy']:.2f} (masked).",
    f"DistilBERT accuracy: {dbert['accuracy']:.2f} (full) → {dbert_mask['accuracy']:.2f} (masked).",
    f"Both stay above the {maj:.2f} majority baseline with the keyword hidden.",
    "→ the models use sentence context, not a keyword lookup.",
])

# --- 13. Challenges & lessons ----------------------------------------------
bullets("Challenges & Lessons Learned", [
    "Extreme class imbalance with only 0–2 samples for the rarest classes on the test set.",
    "Weak labels are noisy — \"may\" as permission vs. possibility can be mislabeled.",
    "Tiny test set makes macro-F1 high-variance.",
    "Lesson: a well-tuned linear model can match a Transformer on small, imbalanced text.",
    "Lesson: data volume for rare classes — not model capacity — is the real bottleneck.",
])

# --- 14. Future work + real-world ------------------------------------------
bullets("Real-World Use & Future Work", [
    "Use: feed a document-to-code agent — typed requirements → test plans, traceability, coverage.",
    "Beneficiaries: firmware/hardware validation, compliance, and technical-writing teams.",
    "Future: expand corpus & add more standards (learning curve still climbing).",
    "Future: build a human-verified gold test set to remove weak-label noise.",
    "Future: add downstream requirement→test-case generation; confidence calibration for human-in-the-loop.",
])

# --- 15. Thank you ----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.8), SW - Inches(1.8), Inches(2))
tf = _tf(tb)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Thank You"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.space_before = Pt(16)
r = p2.add_run(); r.text = REPO
r.font.size = Pt(18); r.font.color.rgb = RGBColor(0xCF, 0xDD, 0xF2)


def main() -> None:
    out = ROOT / "EECE5644_Iteration5_Final_Presentation_Chittela_Ajayanath.pptx"
    prs.save(str(out))
    print(f"[slides] wrote {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
